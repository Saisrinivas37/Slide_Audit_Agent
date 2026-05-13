"""
=============================================================================
  PPTX Slide Audit Agent — JedAI Edition (Cadence)
  -----------------------------------------------------------------------
  Purpose : Compare training PPTX slides against latest product manuals
            (What's New / User Guide / Command Reference) and produce an
            Excel audit report with:
              Sheet 1 – "Slide Audit"   : changes needed in existing slides
              Sheet 2 – "New Commands"  : new commands/params to consider adding
  Inputs  : 1) PPTX file path
            2) One or more manual files (PDF / DOCX / TXT / PPTX)
  Output  : Excel (.xlsx) audit report
=============================================================================
"""

import os
import re
import json
import time
import textwrap
import argparse
import getpass
import requests
from pathlib import Path
from concurrent.futures import ThreadPoolExecutor, as_completed

# ── Optional imports (install if missing) ──────────────────────────────────
try:
    from pptx import Presentation
except ImportError:
    raise ImportError("Install python-pptx:  pip install python-pptx")

try:
    import openpyxl
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
    from openpyxl.utils import get_column_letter
except ImportError:
    raise ImportError("Install openpyxl:  pip install openpyxl")

try:
    import fitz  # PyMuPDF
except ImportError:
    fitz = None
    print("[WARN] PyMuPDF not found. PDF parsing disabled. Install: pip install PyMuPDF")

try:
    from docx import Document as DocxDocument
except ImportError:
    DocxDocument = None
    print("[WARN] python-docx not found. DOCX manual parsing disabled. Install: pip install python-docx")


# ═══════════════════════════════════════════════════════════════════════════
# 1. CONFIGURATION
# ═══════════════════════════════════════════════════════════════════════════
# ── JedAI Endpoints (extracted from your notebook) ──
LOGIN_URL      = "https://jedai-ai.cadence.com:2514/api/v1/security/login"
COMPLETION_URL = "https://jedai-ai.cadence.com:2514/api/copilot/v1/llm/completion"
CHAT_URLS      = [  # Candidate chat endpoints (tried in order for chat-style models)
    "https://jedai-ai.cadence.com:2514/api/copilot/v1/llm/chat/completion",
    "https://jedai-ai.cadence.com:2514/api/copilot/v1/llm/chat/completions",
    "https://jedai-ai.cadence.com:2514/api/copilot/v1/chat/completions",
]
JEDAI_MODEL    = "GCP_claude-opus-4-6"  # Set to "auto" to pick the first available model, a specific name (e.g. "GCP_gemini-2.5-pro"), or None to choose at runtime

# ── Your Cadence LDAP credentials (will be prompted at runtime if not provided via CLI) ──
JEDAI_USERNAME = None
JEDAI_PASSWORD = None

# ── Product Mode ──
# Supported: "innovus", "genus", "conformal", "tempus", "joules", "pegasus", "modus", "generic"
# "generic" uses underscore-based detection. Other modes enable product-specific patterns.
PRODUCT_MODE = "generic"

# ── Tuning parameters ──
MAX_TOKENS      = 8192
CHUNK_SIZE      = 8000    # characters per chunk (smaller = more thorough per-chunk checking)
SLIDE_CHUNK_SIZE = 4000   # characters per slide chunk (for command extraction)
REQUEST_TIMEOUT = 300     # seconds
RETRY_ATTEMPTS  = 3
RETRY_DELAY     = 5       # seconds between retries
MAX_WORKERS     = 5       # max parallel API calls (increase for faster runs, decrease if rate-limited)
ENABLE_PHASE2   = False   # Set to True to scan for new commands in manual (slower), False to skip

# ── Global token (set after login) ──
AUTH_TOKEN = None
AVAILABLE_MODELS = []  # Populated after model discovery


# ═══════════════════════════════════════════════════════════════════════════
# 2. JEDAI AUTHENTICATION (Two-Step: Login → Token)
# ═══════════════════════════════════════════════════════════════════════════

def jedai_login():
    """
    Authenticate with JedAI using LDAP credentials.
    Returns the access_token string.
    """
    global AUTH_TOKEN
    print("  🔐 Authenticating with JedAI (LDAP) ...")

    login_payload = {
        "username": JEDAI_USERNAME,
        "password": JEDAI_PASSWORD,
        "provider": "LDAP"
    }

    for attempt in range(1, RETRY_ATTEMPTS + 1):
        try:
            resp = requests.post(
                LOGIN_URL,
                headers={"Content-Type": "application/json"},
                json=login_payload,
                timeout=30,
                verify=False  # Internal Cadence cert — adjust if needed
            )

            if resp.status_code == 200:
                AUTH_TOKEN = resp.json()["access_token"]
                print(f"  ✅ Login successful! Token: {AUTH_TOKEN[:20]}...")
                return AUTH_TOKEN
            else:
                print(f"  ❌ Login failed (HTTP {resp.status_code}): {resp.text}")
                if attempt < RETRY_ATTEMPTS:
                    time.sleep(RETRY_DELAY)

        except requests.exceptions.RequestException as e:
            print(f"  ⚠️  Login attempt {attempt}/{RETRY_ATTEMPTS} failed: {e}")
            if attempt < RETRY_ATTEMPTS:
                time.sleep(RETRY_DELAY)

    raise RuntimeError("❌ JedAI authentication failed after all retries. Check credentials.")


def prompt_model_selection(reason: str = "") -> None:
    """
    Show the available models list and let the user pick one.
    Updates JEDAI_MODEL globally.
    """
    global JEDAI_MODEL
    if not AVAILABLE_MODELS:
        JEDAI_MODEL = input("  Enter the JedAI model name to use: ").strip()
        if not JEDAI_MODEL:
            raise RuntimeError("No model name provided. Cannot proceed.")
        return

    if reason:
        print(f"\n  {reason}")
    print(f"  📋 Available models:")
    for i, name in enumerate(AVAILABLE_MODELS, start=1):
        print(f"      {i:3d}. {name}")

    while True:
        choice = input(
            f"\n  Enter model number (1-{len(AVAILABLE_MODELS)}) or name "
            f"[press Enter to auto-select '{AVAILABLE_MODELS[0]}']: "
        ).strip()

        if not choice:
            JEDAI_MODEL = AVAILABLE_MODELS[0]
            print(f"  ✅ Auto-selected model: {JEDAI_MODEL}")
            return
        elif choice.isdigit():
            idx = int(choice)
            if 1 <= idx <= len(AVAILABLE_MODELS):
                JEDAI_MODEL = AVAILABLE_MODELS[idx - 1]
                print(f"  ✅ Selected model: {JEDAI_MODEL}")
                return
            else:
                print(f"  ❌ Invalid number. Must be 1-{len(AVAILABLE_MODELS)}. Try again.")
        else:
            if choice in AVAILABLE_MODELS:
                JEDAI_MODEL = choice
                print(f"  ✅ Selected model: {JEDAI_MODEL}")
                return
            else:
                print(f"  ❌ Model '{choice}' not in the list. Try again.")


def discover_models():
    """
    Query JedAI API to list available models.
    If a model was set in the config, validate it against the list.
    Otherwise, prompt the user to pick one (or press Enter to auto-select).
    """
    global JEDAI_MODEL, AVAILABLE_MODELS
    user_provided_model = JEDAI_MODEL  # May have been set in config

    print("  🔍 Discovering available models ...")
    # Try common JedAI model-listing endpoints
    base_url = LOGIN_URL.rsplit('/api/', 1)[0]
    model_urls = [
        f"{base_url}/api/copilot/v1/llm/models",
        f"{base_url}/api/v1/models",
        f"{base_url}/api/copilot/v1/models",
    ]

    for url in model_urls:
        try:
            resp = requests.get(
                url,
                headers={"Authorization": f"Bearer {AUTH_TOKEN}", "accept": "application/json"},
                timeout=30,
                verify=False,
            )
            if resp.status_code == 200:
                data = resp.json()
                models = []
                if isinstance(data, list):
                    models = data
                elif isinstance(data, dict):
                    models = data.get("models", data.get("data", []))

                if models:
                    AVAILABLE_MODELS = []
                    for m in models:
                        name = m if isinstance(m, str) else m.get("id", m.get("name", m.get("model", str(m))))
                        AVAILABLE_MODELS.append(name)

                    # If set to "auto", pick the first available model
                    if user_provided_model and user_provided_model.lower() == "auto":
                        JEDAI_MODEL = AVAILABLE_MODELS[0]
                        print(f"  ✅ Auto-selected model: {JEDAI_MODEL}")
                        return

                    # If user provided a specific model name, validate it
                    if user_provided_model:
                        if user_provided_model in AVAILABLE_MODELS:
                            print(f"  ✅ Using configured model: {user_provided_model}")
                            return
                        else:
                            prompt_model_selection(
                                f"⚠️  Configured model '{user_provided_model}' not found in available models."
                            )
                            return

                    prompt_model_selection()
                    return
        except requests.exceptions.RequestException:
            continue

    # If discovery fails, use user-provided model or prompt
    if user_provided_model:
        print(f"  ⚠️  Could not discover models. Using provided model: {user_provided_model}")
        return
    print("  ⚠️  Could not auto-discover models.")
    JEDAI_MODEL = input("  Enter the JedAI model name to use: ").strip()
    if not JEDAI_MODEL:
        raise RuntimeError("No model name provided. Cannot proceed.")


# ═══════════════════════════════════════════════════════════════════════════
# 3. FILE PARSERS
# ═══════════════════════════════════════════════════════════════════════════

def parse_pptx(filepath: str) -> list:
    """
    Extract text from each slide in a PPTX file.
    Returns: [{"slide_number": int, "text": str}, ...]
    """
    import zipfile

    # Pre-check: test the PPTX zip for corrupt entries and report media per slide
    try:
        with zipfile.ZipFile(filepath, 'r') as z:
            bad_files = z.testzip()
            if bad_files is not None:
                # Find which slides reference the bad media
                print(f"\n  ⚠️  Corrupt file detected in PPTX: {bad_files}")
                bad_basename = bad_files.split("/")[-1]  # e.g. media1.m4a
                # Scan slide relationship files to find which slide uses it
                for name in z.namelist():
                    if name.startswith("ppt/slides/_rels/slide") and name.endswith(".xml.rels"):
                        try:
                            rel_content = z.read(name).decode("utf-8", errors="ignore")
                            if bad_basename in rel_content:
                                # Extract slide number from filename like slide3.xml.rels
                                slide_file = name.split("/")[-1]  # slide3.xml.rels
                                slide_num = re.search(r"slide(\d+)", slide_file)
                                sn = slide_num.group(1) if slide_num else "?"
                                print(f"       → Referenced by slide {sn}")
                        except Exception:
                            pass
                print(f"\n  💡 Fix: Open the PPTX in PowerPoint, remove the media from the")
                print(f"     affected slide(s), then Save As a new file.")
                raise RuntimeError(
                    f"PPTX file is corrupt ({bad_files}). "
                    f"Please repair it in PowerPoint and try again."
                )
    except zipfile.BadZipFile as e:
        raise RuntimeError(f"PPTX file is not a valid zip archive: {e}")

    prs = Presentation(filepath)
    slides_data = []
    for idx, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = para.text.strip()
                    if line:
                        texts.append(line)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        cell_text = cell.text.strip()
                        if cell_text:
                            texts.append(cell_text)
        slides_data.append({
            "slide_number": idx,
            "text": "\n".join(texts)
        })
    return slides_data


def parse_pdf(filepath: str) -> str:
    """Extract text from a PDF file using PyMuPDF."""
    if fitz is None:
        raise RuntimeError("PyMuPDF is required for PDF parsing. pip install PyMuPDF")
    doc = fitz.open(filepath)
    pages = []
    for page in doc:
        pages.append(page.get_text())
    doc.close()
    return "\n\n".join(pages)


def parse_docx(filepath: str) -> str:
    """Extract text from a DOCX file."""
    if DocxDocument is None:
        raise RuntimeError("python-docx is required. pip install python-docx")
    doc = DocxDocument(filepath)
    return "\n".join(para.text for para in doc.paragraphs if para.text.strip())


def parse_txt(filepath: str) -> str:
    """Read plain text file."""
    with open(filepath, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()


def parse_manual(filepath: str) -> str:
    """Auto-detect file type and parse accordingly."""
    ext = Path(filepath).suffix.lower()
    parsers = {
        ".pdf":  parse_pdf,
        ".docx": parse_docx,
        ".txt":  parse_txt,
        ".text": parse_txt,
        ".md":   parse_txt,
        ".pptx": lambda fp: "\n\n".join(
            f"[Slide {s['slide_number']}]\n{s['text']}" for s in parse_pptx(fp)
        ),
    }
    parser = parsers.get(ext)
    if parser is None:
        raise ValueError(f"Unsupported file type: {ext}. Supported: {list(parsers.keys())}")
    print(f"  📄 Parsing {Path(filepath).name} ({ext}) ...")
    return parser(filepath)


# ═══════════════════════════════════════════════════════════════════════════
# 4. TEXT CHUNKER
# ═══════════════════════════════════════════════════════════════════════════

def chunk_text(text: str, chunk_size: int = CHUNK_SIZE) -> list:
    """Split text into chunks, respecting paragraph boundaries where possible."""
    if len(text) <= chunk_size:
        return [text]

    chunks = []
    current = ""
    for paragraph in text.split("\n"):
        if len(current) + len(paragraph) + 1 > chunk_size:
            if current:
                chunks.append(current)
            current = paragraph
        else:
            current = f"{current}\n{paragraph}" if current else paragraph
    if current:
        chunks.append(current)
    return chunks


# ═══════════════════════════════════════════════════════════════════════════
# 5. JEDAI API CALL (Completions-Style with `prompt` field)
# ═══════════════════════════════════════════════════════════════════════════

def _is_chat_model(model_name: str) -> bool:
    """Return True if the model expects a chat-style messages payload."""
    chat_prefixes = ("GCP_", "AzureOpenAI_", "Custom_LLM")
    return any(model_name.startswith(p) for p in chat_prefixes)


def call_jedai(prompt_text: str, max_tokens: int = MAX_TOKENS, temperature: float = 0.1) -> str:
    """
    Send a prompt to JedAI's completion endpoint and return the response text.
    Auto-detects whether to use chat `messages` or raw `prompt` based on the model name.
    For chat models, tries multiple candidate chat endpoints.
    """
    global AUTH_TOKEN

    if AUTH_TOKEN is None:
        jedai_login()

    headers = {
        "accept": "*/*",
        "Content-Type": "application/json",
        "Authorization": f"Bearer {AUTH_TOKEN}",
    }

    is_chat = _is_chat_model(JEDAI_MODEL)

    # Build payload: chat models use `messages`, on-prem use `prompt`
    if is_chat:
        payload = {
            "model": JEDAI_MODEL,
            "max_tokens": max_tokens,
            "temperature": temperature,
            "top_p": 0.9,
            "messages": [
                {"role": "user", "content": prompt_text}
            ],
        }
        # Try chat endpoints first, then fall back to completion endpoint
        urls_to_try = CHAT_URLS + [COMPLETION_URL]
    else:
        payload = {
            "model": JEDAI_MODEL,
            "max_tokens": max_tokens,
            "n": 1,
            "temperature": temperature,
            "top_p": 0.9,
            "prompt": prompt_text,
        }
        urls_to_try = [COMPLETION_URL]

    for attempt in range(1, RETRY_ATTEMPTS + 1):
        last_error = None
        for url in urls_to_try:
            try:
                resp = requests.post(
                    url,
                    headers=headers,
                    json=payload,
                    timeout=REQUEST_TIMEOUT,
                    verify=False  # Internal Cadence cert
                )

                # ── Handle token expiry: re-login and retry ──
                if resp.status_code == 401:
                    print("  🔄 Token expired, re-authenticating ...")
                    jedai_login()
                    headers["Authorization"] = f"Bearer {AUTH_TOKEN}"
                    break  # retry this attempt with refreshed token

                # ── Try next URL if this endpoint doesn't work ──
                if resp.status_code in (400, 404, 405):
                    last_error = f"HTTP {resp.status_code} from {url}: {resp.text[:200]}"
                    continue  # try next URL

                if resp.status_code != 200:
                    print(f"  ⚠️  HTTP {resp.status_code} response: {resp.text[:500]}")

                resp.raise_for_status()
                data = resp.json()

                # ── Extract text from response ──
                if "choices" in data:
                    choice = data["choices"][0]
                    # Handle both completion and chat response formats
                    return choice.get("text", choice.get("message", {}).get("content", str(data)))
                elif "response" in data:
                    return data["response"]
                elif "output" in data:
                    return data["output"]
                else:
                    return json.dumps(data)

            except requests.exceptions.RequestException as e:
                last_error = str(e)
                continue  # try next URL
        else:
            # All URLs exhausted for this attempt
            if last_error:
                print(f"  ⚠️  Attempt {attempt}/{RETRY_ATTEMPTS}: all endpoints failed. Last error: {last_error}")
                if attempt < RETRY_ATTEMPTS:
                    time.sleep(RETRY_DELAY)
                    continue
                # Final attempt failed — let user pick a different model
                print(f"  ❌ Model '{JEDAI_MODEL}' failed on all endpoints after {RETRY_ATTEMPTS} attempts.")
                prompt_model_selection(
                    f"⚠️  Model '{JEDAI_MODEL}' is not compatible. Please choose a different model."
                )
                # Rebuild payload and URL list for the newly selected model
                is_chat = _is_chat_model(JEDAI_MODEL)
                if is_chat:
                    payload = {
                        "model": JEDAI_MODEL,
                        "max_tokens": max_tokens,
                        "temperature": temperature,
                        "top_p": 0.9,
                        "messages": [
                            {"role": "user", "content": prompt_text}
                        ],
                    }
                    urls_to_try = CHAT_URLS + [COMPLETION_URL]
                else:
                    payload = {
                        "model": JEDAI_MODEL,
                        "max_tokens": max_tokens,
                        "n": 1,
                        "temperature": temperature,
                        "top_p": 0.9,
                        "prompt": prompt_text,
                    }
                    urls_to_try = [COMPLETION_URL]
                # Retry from scratch with new model
                return call_jedai(prompt_text, max_tokens, temperature)


# ═══════════════════════════════════════════════════════════════════════════
# 6. COMMAND EXTRACTION (Regex + LLM)
# ═══════════════════════════════════════════════════════════════════════════

# Regex patterns for EDA command syntax
_RE_COMMAND   = re.compile(r'\b([a-zA-Z][a-zA-Z0-9]*(?:_[a-zA-Z0-9]+)+)\b')       # words_with_underscores ONLY
_RE_OPTION    = re.compile(r'(?<!\w)(-[a-zA-Z][a-zA-Z0-9_]*)\b')                   # -option_name
# Attribute patterns: value after set/set_db/get_db, or word near "attribute"
_RE_SET_ATTR  = re.compile(
    r'\b(?:set_db|set|get_db|get)\s+[\[\]$\w./]*\s+'     # set_db [obj] <attr>
    r'\.?([a-zA-Z][a-zA-Z0-9_]*(?:\.[a-zA-Z][a-zA-Z0-9_]*)*)',  # capture attr name
    re.IGNORECASE
)
_RE_DOT_ATTR  = re.compile(r'\.([a-zA-Z][a-zA-Z0-9_]+(?:\.[a-zA-Z][a-zA-Z0-9_]+)*)')  # .attr_name
_RE_ATTR_WORD = re.compile(
    r'\b[Aa]ttribut\w*\s+[:\-]?\s*'                       # "attribute:" or "Attribute -"
    r'\.?([a-zA-Z][a-zA-Z0-9_]+(?:\.[a-zA-Z][a-zA-Z0-9_]+)*)',
    re.IGNORECASE
)

# ── Conformal-specific patterns ──
# Conformal commands are multi-word (space-separated): read design, set system mode,
# report verification, add pin constraints, etc.
# Known Conformal command verbs (first word of multi-word commands)
_CONFORMAL_VERBS = {
    "read", "write", "set", "get", "report", "add", "remove", "delete",
    "analyze", "compare", "map", "unmap", "flatten", "model", "rename",
    "save", "restore", "run", "abort", "check", "clean", "configure",
    "create", "define", "diagnose", "do", "exit", "find", "fix",
    "identify", "ignore", "isolate", "load", "mark", "merge",
    "modify", "print", "put", "replace", "reset", "resolve",
    "source", "specify", "split", "start", "tclmode", "undo",
    "usage", "verify", "vpx",
}

# Conformal multi-word command pattern: verb followed by 1-4 more words
# e.g. "read design", "set system mode", "report verification -summary"
_RE_CONFORMAL_CMD = re.compile(
    r'\b(' + '|'.join(sorted(_CONFORMAL_VERBS)) + r')'
    r'((?:\s+[a-zA-Z][a-zA-Z0-9_]*){1,4})',
    re.IGNORECASE
)

# Conformal "set" command: "set <noun_phrase> <value>" — the noun phrase is the attribute
# e.g. "set system mode lec" → command="set system mode", attribute value="lec"
_RE_CONFORMAL_SET_ATTR = re.compile(
    r'\b(set\s+(?:[a-zA-Z][a-zA-Z0-9_]*\s+){1,3})'    # "set system mode "
    r'([a-zA-Z0-9_/\-\.]+)',                            # the value being set
    re.IGNORECASE
)

# Conformal noise — common English phrases that match the verb pattern but aren't commands
_CONFORMAL_NOISE = {
    "set of", "set the", "set up", "set it", "read the", "read more",
    "read about", "run the", "run it", "add the", "add a", "add an",
    "remove the", "remove a", "check the", "check if", "check for",
    "report the", "report a", "find the", "find a", "get the", "get a",
    "create a", "create the", "define the", "define a",
    "compare the", "compare a", "write the", "write a",
    "save the", "save a", "load the", "load a",
}

# Common EDA noise words to filter out (section titles, generic terms)
_NOISE_WORDS = {
    # Generic words that happen to have underscores but are NOT EDA commands
    "slide_number", "slide_numbers", "slide_audit", "new_commands",
    "change_type", "change_description", "action_needed", "section_chapter",
    "include_in_pptx", "file_name", "file_path", "file_type",
    "click_here", "right_click", "left_click", "double_click",
    "end_of", "table_of", "list_of", "set_of", "number_of",
    "one_of", "each_of", "all_of", "any_of", "none_of",
    "e_g", "i_e", "etc_", "vs_", "re_",
}


def extract_commands_conformal(slides_data: list) -> list:
    """
    Extract Conformal-style multi-word commands, options, and attributes.
    Conformal commands: "read design", "set system mode", "report verification"
    Options: -golden, -revised, -noaliases (same dash-prefix as other tools)
    Attributes: values set via "set" commands (e.g. "set system mode lec" → attr: lec)
    Returns: [{"command": str, "options": [str], "attributes": [str],
               "context": str, "slide_numbers": str}, ...]
    """
    found = {}

    for slide in slides_data:
        snum = str(slide["slide_number"])
        text = slide["text"]
        if not text.strip():
            continue

        lines = text.split("\n")
        for line in lines:
            line_stripped = line.strip()
            if not line_stripped:
                continue

            # ── Find Conformal multi-word commands ──
            commands_in_line = []
            for m in _RE_CONFORMAL_CMD.finditer(line_stripped):
                verb = m.group(1).lower()
                rest = m.group(2).strip()
                # Build the full command (lowercase for consistency)
                full_cmd = f"{verb} {rest}".strip().lower()
                # Filter out noise phrases
                if full_cmd in _CONFORMAL_NOISE:
                    continue
                # Filter out phrases ending in common articles/prepositions
                last_word = full_cmd.split()[-1]
                if last_word in ("the", "a", "an", "of", "to", "in", "on", "for", "is", "are", "it", "if"):
                    continue
                # Must have at least 2 words
                if len(full_cmd.split()) < 2:
                    continue
                commands_in_line.append(full_cmd)

            # ── Also find underscore-style commands (some Conformal uses both) ──
            for m in _RE_COMMAND.finditer(line_stripped):
                cmd = m.group(1)
                if cmd.lower() in _NOISE_WORDS:
                    continue
                if len(cmd) < 4:
                    continue
                if cmd.startswith("__") or cmd.endswith("__"):
                    continue
                commands_in_line.append(cmd)

            # ── Find options/flags (-prefixed) ──
            options_in_line = []
            for m in _RE_OPTION.finditer(line_stripped):
                opt = m.group(1)
                if len(opt) < 3:
                    continue
                options_in_line.append(opt)

            # ── Find attributes (Conformal "set" attributes) ──
            attrs_in_line = set()
            for m in _RE_CONFORMAL_SET_ATTR.finditer(line_stripped):
                value = m.group(2).strip()
                if value and len(value) >= 2 and not value.startswith("-"):
                    attrs_in_line.add(value)

            # Also catch dot-attributes and standard set_db attributes
            for m in _RE_DOT_ATTR.finditer(line_stripped):
                attr = m.group(1).strip()
                if attr and len(attr) >= 3 and attr.lower() not in _NOISE_WORDS:
                    attrs_in_line.add(attr)

            # ── Associate options/attrs with commands ──
            if commands_in_line:
                for cmd in commands_in_line:
                    if cmd not in found:
                        ctx = line_stripped[:150]
                        found[cmd] = {"options": set(), "attributes": set(),
                                      "context": ctx, "slides": set()}
                    found[cmd]["slides"].add(snum)
                    for opt in options_in_line:
                        found[cmd]["options"].add(opt)
                    for attr in attrs_in_line:
                        if attr != cmd:
                            found[cmd]["attributes"].add(attr)
            elif (options_in_line or attrs_in_line) and found:
                last_cmd = None
                for cmd_name in reversed(list(found.keys())):
                    if snum in found[cmd_name]["slides"]:
                        last_cmd = cmd_name
                        break
                if last_cmd:
                    for opt in options_in_line:
                        found[last_cmd]["options"].add(opt)
                    for attr in attrs_in_line:
                        if attr != last_cmd:
                            found[last_cmd]["attributes"].add(attr)

    # Convert to list format
    result = []
    for cmd_name, info in sorted(found.items()):
        slide_nums = ",".join(sorted(info["slides"],
                                      key=lambda x: int(x) if x.isdigit() else 0))
        opts = sorted(info["options"])
        attrs = sorted(info["attributes"])
        result.append({
            "command": cmd_name,
            "options": opts,
            "attributes": attrs,
            "context": info["context"],
            "slide_numbers": slide_nums,
        })
    return result


def extract_commands_regex(slides_data: list) -> list:
    """
    Extract EDA commands along with their associated options and attributes.
    Dispatches to Conformal-specific extraction if PRODUCT_MODE == 'conformal'.
    Groups options and attributes under their parent command.
    Returns: [{"command": str, "options": [str], "attributes": [str],
               "context": str, "slide_numbers": str}, ...]
    """
    # If Conformal mode, use the multi-word command extractor
    if PRODUCT_MODE.lower() == "conformal":
        return extract_commands_conformal(slides_data)

    # {command_name: {"options": set, "attributes": set, "context": str, "slides": set}}
    found = {}

    for slide in slides_data:
        snum = str(slide["slide_number"])
        text = slide["text"]
        if not text.strip():
            continue

        lines = text.split("\n")
        for line in lines:
            line_stripped = line.strip()
            if not line_stripped:
                continue

            # ── Find commands (ONLY words_with_underscores — no single words) ──
            commands_in_line = []
            for m in _RE_COMMAND.finditer(line_stripped):
                cmd = m.group(1)
                if cmd.lower() in _NOISE_WORDS:
                    continue
                if len(cmd) < 4:
                    continue
                if cmd.startswith("__") or cmd.endswith("__"):
                    continue
                commands_in_line.append(cmd)

            # ── Find options/flags (-prefixed) in this line ──
            options_in_line = []
            for m in _RE_OPTION.finditer(line_stripped):
                opt = m.group(1)
                if len(opt) < 3:
                    continue
                options_in_line.append(opt)

            # ── Find attributes ──
            attrs_in_line = set()

            # Pattern 1: set_db / set / get_db followed by attribute name
            for m in _RE_SET_ATTR.finditer(line_stripped):
                attr = m.group(1).strip()
                if attr and len(attr) >= 3 and attr.lower() not in _NOISE_WORDS:
                    attrs_in_line.add(attr)

            # Pattern 2: dot-prefixed attributes (.attr_name)
            for m in _RE_DOT_ATTR.finditer(line_stripped):
                attr = m.group(1).strip()
                if attr and len(attr) >= 3 and attr.lower() not in _NOISE_WORDS:
                    attrs_in_line.add(attr)

            # Pattern 3: word "attribute" followed by name
            for m in _RE_ATTR_WORD.finditer(line_stripped):
                attr = m.group(1).strip()
                if attr and len(attr) >= 3 and attr.lower() not in _NOISE_WORDS:
                    attrs_in_line.add(attr)

            # ── Associate options and attributes with the command on the same line ──
            if commands_in_line:
                for cmd in commands_in_line:
                    if cmd not in found:
                        ctx = line_stripped[:150]
                        found[cmd] = {"options": set(), "attributes": set(),
                                      "context": ctx, "slides": set()}
                    found[cmd]["slides"].add(snum)
                    for opt in options_in_line:
                        found[cmd]["options"].add(opt)
                    for attr in attrs_in_line:
                        # Don't add the command itself as its own attribute
                        if attr != cmd:
                            found[cmd]["attributes"].add(attr)
            # If line has options/attrs but no command, attach to last command on slide
            elif (options_in_line or attrs_in_line) and found:
                last_cmd = None
                for cmd_name in reversed(list(found.keys())):
                    if snum in found[cmd_name]["slides"]:
                        last_cmd = cmd_name
                        break
                if last_cmd:
                    for opt in options_in_line:
                        found[last_cmd]["options"].add(opt)
                    for attr in attrs_in_line:
                        if attr != last_cmd:
                            found[last_cmd]["attributes"].add(attr)

    # Convert to list format
    result = []
    for cmd_name, info in sorted(found.items()):
        slide_nums = ",".join(sorted(info["slides"],
                                      key=lambda x: int(x) if x.isdigit() else 0))
        opts = sorted(info["options"])
        attrs = sorted(info["attributes"])
        result.append({
            "command": cmd_name,
            "options": opts,
            "attributes": attrs,
            "context": info["context"],
            "slide_numbers": slide_nums,
        })
    return result


def _get_conformal_prompt_rules() -> str:
    """Return Conformal-specific extraction rules for the LLM prompt."""
    return textwrap.dedent("""\
IMPORTANT RULES (Conformal / LEC product):
  - COMMANDS are MULTI-WORD (space-separated), e.g.:
    "read design", "set system mode", "report verification",
    "add pin constraints", "write hier compare dofile",
    "analyze datapath", "map key points"
  - Commands typically start with a VERB: read, write, set, get, report,
    add, remove, analyze, compare, map, flatten, diagnose, verify, etc.
  - Some commands may ALSO use underscore format — include those too.
  - OPTIONS start with hyphen: -golden, -revised, -noaliases, -summary
  - ATTRIBUTES are values configured via "set" commands:
    e.g. "set system mode lec" → command is "set system mode",
         the mode value "lec" is the attribute/setting.
    e.g. "set flatten model -seq_constant" → command is "set flatten model"
  - Do NOT list single English words as commands unless they are known
    Conformal commands (e.g. "tclmode", "vpx").
  - Do NOT list options or attributes as standalone commands.
  - Include the FULL multi-word command (e.g. "read design" not just "read").
""")


def _get_generic_prompt_rules() -> str:
    """Return generic/underscore-based extraction rules for the LLM prompt."""
    return textwrap.dedent("""\
IMPORTANT RULES:
  - COMMANDS MUST contain at least one underscore: set_db, report_timing,
    create_clock. Do NOT include single-word commands.
  - Group options AND attributes UNDER their parent command.
  - OPTIONS start with hyphen: -effort, -hold, -max_fanout
  - ATTRIBUTES are values set/configured via set_db, set, get_db:
    e.g. in "set_db [current_design] .max_transition 0.5"
         → command is "set_db", attribute is "max_transition"
  - Do NOT list options or attributes as standalone commands.
""")


def build_extract_commands_prompt(slide_chunk: str, regex_commands: str) -> str:
    """
    Prompt to refine regex-extracted commands (with options and attributes grouped).
    Adapts rules based on PRODUCT_MODE.
    """
    if PRODUCT_MODE.lower() == "conformal":
        product_rules = _get_conformal_prompt_rules()
        task_description = textwrap.dedent(f"""\
### TASK ###
Review the slide content below and the pre-extracted list above.
Your job is to:
1. REMOVE any items that are NOT actual Conformal/LEC EDA commands
   (e.g. section titles, headings, generic English phrases).
2. ADD any Conformal commands that were MISSED. Conformal commands are
   MULTI-WORD (e.g. "read design", "set system mode", "report verification").
   Also include any underscore-style commands if present.
   Scan EVERY slide carefully to ensure no commands are missed.
3. For each command, ensure ALL its options AND attributes from the slides
   are listed together.
4. Identify ATTRIBUTES — for Conformal, these are values set via "set"
   commands. For example: "set system mode lec" has attribute "lec".
   Also look for settings/modes/parameters being configured.
""")
    else:
        product_rules = _get_generic_prompt_rules()
        task_description = textwrap.dedent(f"""\
### TASK ###
Review the slide content below and the pre-extracted list above.
Your job is to:
1. REMOVE any items that are NOT actual EDA commands (e.g. section titles,
   headings, generic English phrases, single words without underscores).
2. ADD any EDA commands that were MISSED — but ONLY commands that contain
   at least one underscore (e.g. read_def, set_db, report_timing).
   Do NOT add single-word commands (e.g. route, place, opt, innovus).
   Scan EVERY slide carefully to ensure no underscore-commands are missed.
3. For each command, ensure ALL its options AND attributes from the slides
   are listed together.
4. Identify ATTRIBUTES — these are values set/get via commands like set_db,
   set, get_db, or mentioned near the word "attribute" in the slides.
   Examples: timing_sense, max_transition, max_fanout, clock_gating_cell
""")

    return textwrap.dedent(f"""\
### ROLE ###
You are an expert EDA / VLSI technical analyst at Cadence Design Systems.

### CONTEXT ###
We have already extracted these commands with their options and attributes:

{regex_commands}

{task_description}
{product_rules}
Return a JSON array of objects. Each object must have:
  "command"       : the command name (NOT an option or attribute)
  "options"       : array of options/flags (strings starting with -)
  "attributes"    : array of attribute names set/configured by this command
  "context"       : one-sentence context from the slide
  "slide_numbers" : comma-separated slide numbers where it appears

Return ONLY valid JSON — no markdown fences, no commentary.

### PPTX SLIDE CONTENT ###
{slide_chunk}

### YOUR JSON RESPONSE ###
""")


def build_slide_audit_prompt(command_name: str, command_options: str, command_attrs: str, command_context: str, slide_numbers: str, manual_text: str) -> str:
    """
    Phase 1B prompt: Deep-verify ONE command (with its options and attributes).
    """
    return textwrap.dedent(f"""\
### ROLE ###
You are an expert EDA / VLSI technical auditor at Cadence Design Systems.
Your task is to do a THOROUGH and ACCURATE check of ONE specific command
and ALL its options, parameters, and attributes from the training slides
against the latest product manual.

### COMMAND TO VERIFY ###
Command: {command_name}
Options/flags used in slides: {command_options}
Attributes used in slides: {command_attrs}
Appears on slide(s): {slide_numbers}
Context from slides: {command_context}

### YOUR CHECKLIST — go through EACH item carefully ###
1. Does this command appear in the manual content below?
2. Has the command been RENAMED to something else?
3. Has the command been DEPRECATED or REMOVED?
   (ONLY if the manual EXPLICITLY says so — do NOT assume removal
    just because the command is not mentioned or only partially covered.)
4. Has the command SYNTAX changed? (different flags, arguments, argument order)
5. For EACH option listed above, check:
   a. Has this option been renamed, deprecated, or removed?
      (ONLY if EXPLICITLY stated in the manual — absence does NOT mean removed.)
   b. Have the valid values for this option changed?
   c. Has the default value for this option changed?
   d. Has the behavior of this option changed?
6. For EACH attribute listed above, check:
   a. Has this attribute been renamed, deprecated, or removed?
      (ONLY if EXPLICITLY stated in the manual — absence does NOT mean removed.)
   b. Have the valid values or value range changed?
   c. Has the default value changed?
   d. Has the behavior or scope changed?
7. Are there any NEW options or attributes added to this command?
8. Have any CONSTRAINTS, LIMITATIONS, or PREREQUISITES changed?

### RESPONSE FORMAT ###
Return a JSON object with exactly one key:

"slide_audit" — an array of objects. Include one entry for EACH change found.
If NOTHING changed or the command is not mentioned in this manual content,
return an empty array.

Each entry must have:
  "command"             : "{command_name}"
  "change_type"         : MUST clearly indicate WHAT changed. Use one of:
                          "Command Renamed", "Command Deprecated", "Command Removed",
                          "Command Syntax Changed", "Command Behavior Changed",
                          "Option Added", "Option Removed", "Option Renamed",
                          "Option Default Changed", "Option Behavior Changed",
                          "Attribute Added", "Attribute Removed", "Attribute Renamed",
                          "Attribute Default Changed", "Attribute Value Changed"
  "change_description"  : SPECIFIC description — include the affected option name
                          (e.g. "-effort") or attribute name (e.g. "max_transition")
                          and quote old vs new values/syntax.
  "slide_numbers"       : "{slide_numbers}"
  "priority"            : "High" (deprecated/removed/renamed),
                          "Medium" (syntax/default/options changed),
                          or "Low" (minor behavior change)
  "action_needed"       : exactly what to change in the slide

IMPORTANT:
  - The "change_type" MUST start with "Command", "Option", or "Attribute"
    to clearly indicate the category of change.
  - For option changes, include the option name in "change_description"
    (e.g. "Option -hold: default changed from 0.1 to 0.05")
  - For attribute changes, include the attribute name in "change_description"
    (e.g. "Attribute max_transition: value range changed from 0-1.0 to 0-2.0")
  - Report EACH change as a SEPARATE entry (don't combine multiple changes).
  - Be PRECISE. Quote the exact old value from slides and new value from manual.
  - Do NOT hallucinate changes. Only report what you can verify from the text.
  - If the command is NOT mentioned in the manual text below, return empty array.
  - CRITICAL: Do NOT assume a command, option, or attribute has been "Removed",
    "Deprecated", or "Deleted" simply because it is not mentioned in the manual
    text. The manual content provided may be partial (e.g. a "What's New" guide,
    a single chapter, or an excerpt). Absence from the manual does NOT mean
    removal. Only report removal or deprecation if the manual EXPLICITLY states
    that the item has been removed, deprecated, or replaced.
  - Similarly, do NOT report "Option Removed" or "Attribute Removed" just
    because a specific option or attribute is not listed in the manual section.
    Manuals often do not exhaustively list every option. Only report removal
    if there is EXPLICIT evidence (e.g. "removed in version X", "no longer
    supported", "deprecated").

Return ONLY valid JSON — no markdown fences, no commentary.

### MANUAL CONTENT ###
{manual_text}

### YOUR JSON RESPONSE ###
""")


def build_new_commands_prompt(slide_commands_summary: str, manual_chunk: str, chunk_idx: int, total_chunks: int) -> str:
    """
    Phase 2 prompt (SECONDARY): Quickly identify new commands in the manual
    that are NOT already covered in the slides.
    """
    return textwrap.dedent(f"""\
### ROLE ###
You are an EDA technical analyst at Cadence Design Systems.

### TASK ###
Scan the manual content below and identify commands, attributes, or parameters
that are NEW (not present in the slide commands list).
This is a secondary/quick scan — only report clearly NEW items worth adding
to training slides. Skip trivial or internal-only items.

Return a JSON object with exactly one key:

"new_commands" — an array of objects, each with:
  "command"         : the new command / attribute / parameter name
  "section_chapter" : section or chapter in the manual where it appears
  "description"     : concise description of what it does
  "include_in_pptx" : "Yes", "Maybe", or "No"
  "rationale"       : brief reason

If nothing new is found, return an empty array.
Return ONLY valid JSON — no markdown fences, no commentary.

### EXISTING SLIDE COMMANDS (already covered) ###
{slide_commands_summary}

### MANUAL CONTENT (chunk {chunk_idx}/{total_chunks}) ###
{manual_chunk}

### YOUR JSON RESPONSE ###
""")


# ═══════════════════════════════════════════════════════════════════════════
# 7. RESPONSE PARSER
# ═══════════════════════════════════════════════════════════════════════════

def parse_jedai_response(raw: str) -> dict:
    """Extract JSON from JedAI response, handling markdown fences if present."""
    # Strip markdown code fences if model wraps the JSON
    cleaned = re.sub(r"```(?:json)?", "", raw).strip().rstrip("`")

    try:
        return json.loads(cleaned)
    except json.JSONDecodeError:
        # Try to find JSON object within the text
        match = re.search(r"\{.*\}", cleaned, re.DOTALL)
        if match:
            try:
                return json.loads(match.group())
            except json.JSONDecodeError:
                pass
    print("  ⚠️  Could not parse JSON from JedAI response. Saving raw for debug.")
    return {"slide_audit": [], "new_commands": [], "_raw_debug": raw[:500]}


# ═══════════════════════════════════════════════════════════════════════════
# 8. EXCEL REPORT WRITER
# ═══════════════════════════════════════════════════════════════════════════

def write_excel_report(slide_audit: list, new_commands: list, output_path: str):
    """Write audit results to a styled Excel workbook with two sheets."""
    wb = openpyxl.Workbook()

    # ── Styles ──
    header_font = Font(name="Calibri", bold=True, color="FFFFFF", size=11)
    header_fill = PatternFill(start_color="2F5496", end_color="2F5496", fill_type="solid")
    high_fill   = PatternFill(start_color="FFC7CE", end_color="FFC7CE", fill_type="solid")
    med_fill    = PatternFill(start_color="FFEB9C", end_color="FFEB9C", fill_type="solid")
    low_fill    = PatternFill(start_color="C6EFCE", end_color="C6EFCE", fill_type="solid")
    thin_border = Border(
        left=Side(style="thin"), right=Side(style="thin"),
        top=Side(style="thin"),  bottom=Side(style="thin"),
    )
    wrap_align  = Alignment(wrap_text=True, vertical="top")

    def style_header(ws, num_cols):
        for col in range(1, num_cols + 1):
            cell = ws.cell(row=1, column=col)
            cell.font = header_font
            cell.fill = header_fill
            cell.alignment = Alignment(horizontal="center", vertical="center")
            cell.border = thin_border

    def auto_width(ws, num_cols, max_width=50):
        for col in range(1, num_cols + 1):
            max_len = 0
            for row in ws.iter_rows(min_col=col, max_col=col, values_only=False):
                for cell in row:
                    if cell.value:
                        max_len = max(max_len, len(str(cell.value)))
            ws.column_dimensions[get_column_letter(col)].width = min(max_len + 4, max_width)

    # ────────────────────── Sheet 1: Slide Audit ──────────────────────────
    ws1 = wb.active
    ws1.title = "Slide Audit"
    headers_1 = ["#", "Command", "Change Type", "Change Description",
                  "Slide Number(s)", "Priority", "Action Needed"]
    ws1.append(headers_1)
    style_header(ws1, len(headers_1))

    # Color fills for change categories
    cmd_change_fill  = PatternFill(start_color="D6DCE4", end_color="D6DCE4", fill_type="solid")   # grey-blue for command changes
    opt_change_fill  = PatternFill(start_color="BDD7EE", end_color="BDD7EE", fill_type="solid")   # light blue for option changes
    attr_change_fill = PatternFill(start_color="E2EFDA", end_color="E2EFDA", fill_type="solid")   # light green for attribute changes
    no_change_fill   = PatternFill(start_color="FFFFFF", end_color="FFFFFF", fill_type="solid")    # white for no change

    # Show ALL commands — merge multiple changes for the same command into one row
    # Step 1: Group entries by command name
    from collections import OrderedDict
    merged = OrderedDict()
    for item in slide_audit:
        cmd = item.get("command", "")
        if cmd not in merged:
            merged[cmd] = {
                "command": cmd,
                "change_types": [],
                "change_descriptions": [],
                "slide_numbers": set(),
                "priorities": [],
                "actions": [],
            }
        # Dedup individual changes
        desc = item.get("change_description", "").strip()
        ctype = item.get("change_type", "").strip()
        dedup_key = (ctype, desc[:80])
        existing_keys = list(zip(merged[cmd]["change_types"],
                                  [d[:80] for d in merged[cmd]["change_descriptions"]]))
        if dedup_key in existing_keys:
            continue
        merged[cmd]["change_types"].append(ctype)
        merged[cmd]["change_descriptions"].append(desc)
        merged[cmd]["priorities"].append(item.get("priority", "Medium"))
        merged[cmd]["actions"].append(item.get("action_needed", "").strip())
        for s in str(item.get("slide_numbers", "")).split(","):
            s = s.strip()
            if s:
                merged[cmd]["slide_numbers"].add(s)

    # Step 2: Write merged rows
    row_idx = 0
    for cmd, info in merged.items():
        row_idx += 1

        # Combine change types, descriptions, and actions with numbered bullets
        if len(info["change_types"]) == 1:
            change_type = info["change_types"][0]
            change_desc = info["change_descriptions"][0]
            action = info["actions"][0]
        else:
            change_type = "; ".join(dict.fromkeys(info["change_types"]))  # unique, preserving order
            change_desc = "\n".join(f"{i}. {d}" for i, d in enumerate(info["change_descriptions"], 1))
            action = "\n".join(f"{i}. {a}" for i, a in enumerate(info["actions"], 1) if a)

        # Highest priority wins
        priority_rank = {"High": 3, "Medium": 2, "Low": 1}
        priority = max(info["priorities"], key=lambda p: priority_rank.get(p, 0))

        slide_nums = ",".join(sorted(info["slide_numbers"],
                                      key=lambda x: int(x) if x.isdigit() else 0))
        row_data = [
            row_idx,
            cmd,
            change_type,
            change_desc,
            slide_nums,
            priority,
            action,
        ]
        ws1.append(row_data)
        row_num = ws1.max_row

        # Priority fill on priority column
        pri_fill = {"High": high_fill, "Medium": med_fill, "Low": low_fill}.get(priority, med_fill)
        ws1.cell(row=row_num, column=6).fill = pri_fill

        # Category fill on change_type column based on prefix
        # For merged rows with multiple change types, pick the most significant category
        if "Command" in change_type:
            cat_fill = cmd_change_fill
        elif "Option" in change_type:
            cat_fill = opt_change_fill
        elif "Attribute" in change_type:
            cat_fill = attr_change_fill
        else:
            cat_fill = no_change_fill
        ws1.cell(row=row_num, column=3).fill = cat_fill

        for col in range(1, len(headers_1) + 1):
            cell = ws1.cell(row=row_num, column=col)
            cell.border = thin_border
            cell.alignment = wrap_align

    auto_width(ws1, len(headers_1))
    ws1.auto_filter.ref = ws1.dimensions

    # ────────────────────── Sheet 2: New Commands ─────────────────────────
    ws2 = wb.create_sheet(title="New Commands")
    headers_2 = ["Command", "Section / Chapter", "Description",
                  "Include in PPTX?", "Rationale"]
    ws2.append(headers_2)
    style_header(ws2, len(headers_2))

    seen_new = set()
    for item in new_commands:
        cmd = item.get("command", "")
        if cmd in seen_new:
            continue
        seen_new.add(cmd)
        row_data = [
            cmd,
            item.get("section_chapter", ""),
            item.get("description", ""),
            item.get("include_in_pptx", "Maybe"),
            item.get("rationale", ""),
        ]
        ws2.append(row_data)
        row_num = ws2.max_row
        for col in range(1, len(headers_2) + 1):
            cell = ws2.cell(row=row_num, column=col)
            cell.border = thin_border
            cell.alignment = wrap_align
        include_val = item.get("include_in_pptx", "Maybe")
        inc_fill = {"Yes": low_fill, "No": high_fill, "Maybe": med_fill}.get(include_val, med_fill)
        ws2.cell(row=row_num, column=4).fill = inc_fill

    auto_width(ws2, len(headers_2))
    ws2.auto_filter.ref = ws2.dimensions

    # ── Save ──
    wb.save(output_path)
    print(f"\n  ✅ Report saved → {output_path}")
    print(f"     Sheet 1 'Slide Audit'  : {row_idx} entries")
    print(f"     Sheet 2 'New Commands'  : {len(seen_new)} entries")


# ═══════════════════════════════════════════════════════════════════════════
# 9. MAIN AGENT ORCHESTRATOR
# ═══════════════════════════════════════════════════════════════════════════

def run_audit(pptx_path: str, manual_paths: list, output_path: str):
    """
    Main pipeline:
      0. Authenticate with JedAI
      1. Parse PPTX → slide text
      2. Parse manual(s) → combined text
      3. Phase 1A: Extract ALL commands/attributes/params from slides (LLM)
      4. Phase 1B: Verify EACH extracted command against manual chunks
      5. Phase 2:  Quick scan for new commands (secondary)
      6. Write Excel report
    """
    print("=" * 65)
    print("  📊  PPTX SLIDE AUDIT AGENT  (JedAI-Powered)")
    print("=" * 65)

    # ── Step 0: Authenticate ──
    print(f"\n[0/6] Authenticating with JedAI ...")
    jedai_login()
    discover_models()

    # ── Step 1: Parse PPTX ──
    print(f"\n[1/6] Parsing PPTX: {pptx_path}")
    slides = parse_pptx(pptx_path)
    slide_text = "\n\n".join(
        f"[Slide {s['slide_number']}]\n{s['text']}" for s in slides if s["text"].strip()
    )
    print(f"  ✔ Extracted text from {len(slides)} slides "
          f"({len(slide_text):,} chars)")

    # ── Step 2: Parse Manual(s) ──
    print(f"\n[2/6] Parsing {len(manual_paths)} manual file(s):")
    manual_texts = []
    for mp in manual_paths:
        txt = parse_manual(mp)
        manual_texts.append(f"--- FILE: {Path(mp).name} ---\n{txt}")
        print(f"  ✔ {Path(mp).name}: {len(txt):,} chars")
    combined_manual = "\n\n".join(manual_texts)

    # ── Step 3 (Phase 1A): Extract commands from PPTX ──
    print(f"\n[3/6] Extracting commands/attributes/params from PPTX slides ...")

    # Step 3a: Regex-based extraction (fast, pattern-based)
    print(f"  🔍 Regex scan for commands, options, and attributes ...")
    regex_extracted = extract_commands_regex(slides)
    print(f"  ✔ Regex found {len(regex_extracted)} commands (with grouped options & attributes)")

    # Step 3b: LLM refinement — validate regex results + catch missed items
    print(f"  🤖 LLM refinement: validating and supplementing regex results ...")
    slide_chunks = chunk_text(slide_text, SLIDE_CHUNK_SIZE)

    # Format regex results as a summary for the LLM prompt
    regex_summary = "\n".join(
        f"  - {item['command']} opts:[{' '.join(item.get('options', []))}] "
        f"attrs:[{' '.join(item.get('attributes', []))}] [slides: {item['slide_numbers']}]"
        for item in regex_extracted
    )

    def _refine_slide_chunk(args):
        """Refine a single slide chunk via LLM (thread-safe)."""
        i, s_chunk, total = args
        prompt = build_extract_commands_prompt(s_chunk, regex_summary)
        raw = call_jedai(prompt)
        parsed = parse_jedai_response(raw)
        if isinstance(parsed, list):
            items = parsed
        elif isinstance(parsed, dict):
            items = parsed.get("commands", parsed.get("data", []))
            if not items:
                for v in parsed.values():
                    if isinstance(v, list):
                        items = v
                        break
        else:
            items = []
        print(f"  ✔ Slide chunk {i}/{total} → {len(items)} items")
        return items

    all_extracted = []
    print(f"  🔄 Refining {len(slide_chunks)} slide chunk(s) in parallel (max {MAX_WORKERS} workers) ...")
    with ThreadPoolExecutor(max_workers=MAX_WORKERS) as executor:
        futures = [executor.submit(_refine_slide_chunk, (i, chunk, len(slide_chunks)))
                   for i, chunk in enumerate(slide_chunks, start=1)]
        for future in as_completed(futures):
            all_extracted.extend(future.result())

    # Merge LLM-refined results with regex results (LLM takes priority for context)
    # Start with regex results as baseline
    seen_cmds = {}
    for item in regex_extracted:
        cmd = item.get("command", "").strip()
        if cmd:
            seen_cmds[cmd] = item

    # Override/add with LLM-refined results
    for item in all_extracted:
        cmd = item.get("command", "").strip()
        if not cmd:
            continue
        # Skip if an option was returned as a standalone command
        if cmd.startswith("-"):
            continue
        # Skip single-word commands (must contain underscore)
        if "_" not in cmd:
            continue
        if cmd in seen_cmds:
            # Merge slide numbers
            old_sn = str(seen_cmds[cmd].get("slide_numbers", ""))
            new_sn = str(item.get("slide_numbers", ""))
            merged = set(s.strip() for s in f"{old_sn},{new_sn}".split(",") if s.strip())
            seen_cmds[cmd]["slide_numbers"] = ",".join(sorted(merged, key=lambda x: int(x) if x.isdigit() else 0))
            if item.get("context"):
                seen_cmds[cmd]["context"] = item["context"]
            # Merge options
            existing_opts = set(seen_cmds[cmd].get("options", []))
            new_opts = set(item.get("options", []))
            seen_cmds[cmd]["options"] = sorted(existing_opts | new_opts)
            # Merge attributes
            existing_attrs = set(seen_cmds[cmd].get("attributes", []))
            new_attrs = set(item.get("attributes", []))
            seen_cmds[cmd]["attributes"] = sorted(existing_attrs | new_attrs)
        else:
            seen_cmds[cmd] = item
    extracted_commands = list(seen_cmds.values())

    print(f"\n  📋 Total unique commands extracted: {len(extracted_commands)}")
    for ec in extracted_commands[:20]:  # Show first 20
        opts = ", ".join(ec.get("options", []))
        attrs = ", ".join(ec.get("attributes", []))
        opts_str = f" opts:[{opts}]" if opts else ""
        attrs_str = f" attrs:[{attrs}]" if attrs else ""
        print(f"      • {ec.get('command', '?')}{opts_str}{attrs_str} (slides: {ec.get('slide_numbers', '?')})")
    if len(extracted_commands) > 20:
        print(f"      ... and {len(extracted_commands) - 20} more")

    # Build a summary of just command names for new-commands prompt
    commands_names_summary = ", ".join(item.get("command", "") for item in extracted_commands)

    # ── Step 4 (Phase 1B): Verify EACH command against full manual ──
    manual_chunks = chunk_text(combined_manual, CHUNK_SIZE)
    total_cmds = len(extracted_commands)
    print(f"\n[4/6] PHASE 1 — Deep-verifying {total_cmds} commands "
          f"against manual ({len(manual_chunks)} chunk(s)) ...")
    print(f"       Each command will be checked against ALL relevant manual sections.")

    all_slide_audit = []
    skipped_commands = []  # Track commands not found in manual

    def _find_relevant_chunks(cmd_name, cmd_options, cmd_attrs):
        """
        Find manual chunks relevant to a command using multiple search strategies:
          1. Search by full command name
          2. Search by command name parts (split on underscore)
          3. Search by option names (without leading dash)
          4. Search by attribute names
        Returns deduplicated list of relevant chunks.
        """
        relevant_set = []  # Use list to preserve order, deduplicate by index
        seen_indices = set()

        # Strategy 1: Full command name
        search_term = cmd_name.lstrip("-.").split()[0].lower()
        for idx, m_chunk in enumerate(manual_chunks):
            if search_term in m_chunk.lower():
                if idx not in seen_indices:
                    relevant_set.append(m_chunk)
                    seen_indices.add(idx)

        # Strategy 2: Command name parts (e.g. "set_db" → search for "set_db",
        # but also for cases where manual uses slightly different formatting)
        # Try without common prefixes like "get_", "set_", "report_"
        parts = cmd_name.split("_")
        if len(parts) > 2:
            # Try the core part (e.g. "create_clock_constraint" → "clock_constraint")
            core = "_".join(parts[1:]).lower()
            if len(core) >= 4:
                for idx, m_chunk in enumerate(manual_chunks):
                    if idx not in seen_indices and core in m_chunk.lower():
                        relevant_set.append(m_chunk)
                        seen_indices.add(idx)

        # Strategy 3: Search by significant options (those with 4+ chars after dash)
        for opt in cmd_options:
            opt_term = opt.lstrip("-").lower()
            if len(opt_term) >= 5:  # Only search for meaningful option names
                for idx, m_chunk in enumerate(manual_chunks):
                    if idx not in seen_indices and opt_term in m_chunk.lower():
                        relevant_set.append(m_chunk)
                        seen_indices.add(idx)

        # Strategy 4: Search by attribute names
        for attr in cmd_attrs:
            attr_term = attr.lower()
            if len(attr_term) >= 5:  # Only search for meaningful attribute names
                for idx, m_chunk in enumerate(manual_chunks):
                    if idx not in seen_indices and attr_term in m_chunk.lower():
                        relevant_set.append(m_chunk)
                        seen_indices.add(idx)

        return relevant_set

    def _verify_command(args):
        """Verify a single command against manual chunks (thread-safe)."""
        cmd_idx, cmd_item, total_cmds = args
        cmd_name = cmd_item.get("command", "")
        cmd_options = cmd_item.get("options", [])
        cmd_attrs = cmd_item.get("attributes", [])
        cmd_context = cmd_item.get("context", "")
        cmd_slides = cmd_item.get("slide_numbers", "")
        options_str = ", ".join(cmd_options) if cmd_options else "(none found in slides)"
        attrs_str = ", ".join(cmd_attrs) if cmd_attrs else "(none found in slides)"

        # Find manual chunks using multiple search strategies
        relevant_chunks = _find_relevant_chunks(cmd_name, cmd_options, cmd_attrs)

        if not relevant_chunks:
            print(f"  ⚠️  [{cmd_idx}/{total_cmds}] {cmd_name} — NOT found in manual "
                  f"(no matching sections for command, options, or attributes)")
            skipped_commands.append(cmd_item)
            return []

        # Combine relevant chunks (up to reasonable size)
        combined_relevant = "\n\n---\n\n".join(relevant_chunks)
        if len(combined_relevant) > CHUNK_SIZE * 3:
            combined_relevant = combined_relevant[:CHUNK_SIZE * 3]

        print(f"  🔍 [{cmd_idx}/{total_cmds}] Verifying: {cmd_name} "
              f"({len(cmd_options)} opts, {len(cmd_attrs)} attrs, slides: {cmd_slides}, "
              f"{len(relevant_chunks)} relevant section(s)) ...")

        prompt = build_slide_audit_prompt(cmd_name, options_str, attrs_str, cmd_context, cmd_slides, combined_relevant)
        raw_response = call_jedai(prompt)
        result = parse_jedai_response(raw_response)

        audit_items = result.get("slide_audit", [])
        if audit_items:
            for ai in audit_items:
                print(f"    ⚠️  [{cmd_name}] CHANGE: {ai.get('change_type', '?')} — {ai.get('change_description', '')[:80]}")
            return audit_items
        else:
            print(f"    ✔ [{cmd_name}] No changes found")
            return []

    print(f"  🚀 Verifying commands in parallel (max {MAX_WORKERS} workers) ...")
    with ThreadPoolExecutor(max_workers=MAX_WORKERS) as executor:
        futures = [executor.submit(_verify_command, (i, cmd_item, total_cmds))
                   for i, cmd_item in enumerate(extracted_commands, start=1)]
        for future in as_completed(futures):
            all_slide_audit.extend(future.result())

    # Report verification coverage
    verified_count = total_cmds - len(skipped_commands)
    print(f"\n  📋 Phase 1 total: {len(all_slide_audit)} slide audit entries")
    print(f"  📊 Verification coverage: {verified_count}/{total_cmds} commands verified "
          f"({100 * verified_count // total_cmds if total_cmds else 0}%)")
    if skipped_commands:
        print(f"  ⚠️  {len(skipped_commands)} command(s) NOT found in any manual section:")
        for sc in skipped_commands:
            opts = ", ".join(sc.get("options", []))
            attrs = ", ".join(sc.get("attributes", []))
            print(f"      • {sc.get('command', '?')} (slides: {sc.get('slide_numbers', '?')})"
                  f"{f' opts: [{opts}]' if opts else ''}"
                  f"{f' attrs: [{attrs}]' if attrs else ''}")
        print(f"  💡 These commands may not be covered by the provided manual(s). "
              f"Consider providing additional manual files.")

    # ── Step 5 (Phase 2): Quick scan for new commands (secondary) ──
    all_new_commands = []
    if ENABLE_PHASE2:
        print(f"\n[5/6] PHASE 2 — Quick scan for new commands (secondary) ...")

        def _scan_manual_chunk(args):
            """Scan a single manual chunk for new commands (thread-safe)."""
            i, m_chunk, total = args
            prompt = build_new_commands_prompt(commands_names_summary, m_chunk, i, total)
            raw_response = call_jedai(prompt)
            result = parse_jedai_response(raw_response)
            new_items = result.get("new_commands", [])
            print(f"  ✔ Manual chunk {i}/{total} → {len(new_items)} new command(s)")
            return new_items

        print(f"  🔄 Scanning {len(manual_chunks)} manual chunk(s) in parallel (max {MAX_WORKERS} workers) ...")
        with ThreadPoolExecutor(max_workers=MAX_WORKERS) as executor:
            futures = [executor.submit(_scan_manual_chunk, (i, chunk, len(manual_chunks)))
                       for i, chunk in enumerate(manual_chunks, start=1)]
            for future in as_completed(futures):
                all_new_commands.extend(future.result())

        print(f"\n  📋 Phase 2 total: {len(all_new_commands)} new commands")
    else:
        print(f"\n[5/6] PHASE 2 — Skipped (ENABLE_PHASE2 = False)")

    # ── Step 6: Write Excel ──
    print(f"\n[6/6] Writing Excel report ...")
    write_excel_report(all_slide_audit, all_new_commands, output_path)
    print("\n" + "=" * 65)
    print("  🎉  AUDIT COMPLETE!")
    print("=" * 65)


# ═══════════════════════════════════════════════════════════════════════════
# 10. CLI ENTRY POINT
# ═══════════════════════════════════════════════════════════════════════════

if __name__ == "__main__":
    parser = argparse.ArgumentParser(
        description="PPTX Slide Audit Agent — Compare slides against latest product manuals using JedAI",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog=textwrap.dedent("""\
        Examples:
          python slide_audit_agent.py --pptx lecture.pptx --manuals whats_new_26.1.pdf cmd_ref.pdf -o audit.xlsx
          python slide_audit_agent.py --pptx deck.pptx --manuals user_guide.docx --output report.xlsx
        """),
    )
    parser.add_argument("--pptx", required=True, help="Path to the PPTX file to audit")
    parser.add_argument("--manuals", nargs="+", required=True, help="Path(s) to manual file(s): PDF, DOCX, TXT, or PPTX")
    parser.add_argument("--product", default="generic",
                        choices=["generic", "innovus", "genus", "conformal", "tempus", "joules", "pegasus", "modus"],
                        help="Product mode — 'conformal' enables multi-word command detection (default: generic)")
    parser.add_argument("-o", "--output", default="slide_audit_report.xlsx", help="Output Excel file path (default: slide_audit_report.xlsx)")
    parser.add_argument("-u", "--username", default=None, help="Cadence LDAP username (overrides script default)")
    parser.add_argument("-p", "--password", default=None, help="Cadence LDAP password (overrides script default)")
    args = parser.parse_args()

    # Set product mode
    PRODUCT_MODE = args.product

    # Set credentials from CLI args or prompt interactively
    if args.username:
        JEDAI_USERNAME = args.username
    else:
        JEDAI_USERNAME = input("Enter your Cadence LDAP username: ").strip()
    if args.password:
        JEDAI_PASSWORD = args.password
    else:
        JEDAI_PASSWORD = getpass.getpass("Enter your Cadence LDAP password: ")

    # Validate files exist
    for fp in [args.pptx] + args.manuals:
        if not os.path.isfile(fp):
            raise FileNotFoundError(f"File not found: {fp}")

    run_audit(args.pptx, args.manuals, args.output)
